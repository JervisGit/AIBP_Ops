"""
Generate: AI Blueprint AO Operations - Design & Approach Validation.pptx
Uses the existing Security & Identity PPTX as a template (preserving IRAS theme/background),
then populates it with Operations content.
"""

import copy
import shutil
import os
import io
import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Constants ──────────────────────────────────────────────────────────────────
SRC = "AI Blueprint AO Security & Identity - Design & Approach Validation.pptx"
DST = "AI Blueprint AO Operations - Design & Approach Validation.pptx"

BLUE   = RGBColor(0x00, 0x5A, 0xAB)   # IRAS blue
RED    = RGBColor(0xFF, 0x00, 0x00)   # Confirmation / Inputs red
BLACK  = RGBColor(0x00, 0x00, 0x00)
GREY   = RGBColor(0x60, 0x60, 0x60)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)
DARK_GREY  = RGBColor(0x40, 0x40, 0x40)

W  = Inches(13.333)
H  = Inches(7.5)

# Tag colour per type
TAG_COLOURS = {
    "For Information":  BLUE,
    "For Confirmation": RED,
    "For Inputs":       RED,
}


# ── Helpers ────────────────────────────────────────────────────────────────────

def _repack_zip(path):
    """
    Rewrite the PPTX zip keeping only the LAST occurrence of each filename.
    This removes stale slide entries left over from the source template,
    preventing PowerPoint from finding the old Security deck slides on repair.
    """
    entries = {}
    with zipfile.ZipFile(path, 'r') as zin:
        for name in zin.namelist():
            entries[name] = zin.read(name)   # last entry wins for duplicates

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)

    with open(path, 'wb') as f:
        f.write(buf.getvalue())


def _remove_all_slides(prs):
    """Remove all slide references from the presentation XML (OPC parts remain but are overwritten)."""
    xml_slides = prs.slides._sldIdLst
    for el in list(xml_slides):
        xml_slides.remove(el)


def _layout(prs, name="Content slide 1"):
    """Return a slide layout by name."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[11]   # fallback


def _add_slide(prs, layout_name="Content slide 1"):
    return prs.slides.add_slide(_layout(prs, layout_name))


def _txbox(slide, text, left, top, width, height,
           font_name="Source Sans Pro", font_pt=12,
           bold=False, color=None, italic=False,
           align=PP_ALIGN.LEFT, wrap=True, para_spacing_before=0):
    """Add a text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # We may receive multi-paragraph text separated by "\n"
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if para_spacing_before and idx > 0:
            p.space_before = Pt(para_spacing_before)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
    return txBox


def _title(slide, text, top=Inches(0.18)):
    """Add the main slide title using the content placeholder (or a text box)."""
    # Use placeholder index 0 if available, else use a text box
    if slide.placeholders:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = text
                # style
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Source Sans Pro"
                        run.font.size = Pt(20)
                return ph
    # Fallback: text box
    return _txbox(slide, text, Inches(0.91), top, Inches(11.5), Inches(0.5),
                  font_pt=20, bold=False, color=BLACK)


def _tag(slide, tag_type):
    """Add the 'For Information / For Confirmation / For Inputs' tag top-right."""
    colour = TAG_COLOURS.get(tag_type, BLUE)
    bold = (tag_type != "For Information")
    _txbox(slide, tag_type,
           left=Inches(10.3), top=Inches(0.12),
           width=Inches(2.9), height=Inches(0.38),
           font_pt=12, bold=bold, color=colour,
           align=PP_ALIGN.RIGHT)


def _section_label(slide, text, top=Inches(0.6)):
    """Coloured section sub-heading just below the title."""
    _txbox(slide, text, Inches(0.91), top, Inches(11.5), Inches(0.35),
           font_pt=11, bold=False, color=GREY)


def _body(slide, text, top, height=Inches(5.5), font_pt=11, bold=False, color=None,
          left=Inches(0.91), width=Inches(11.5)):
    """Generic body text box."""
    _txbox(slide, text, left, top, width, height,
           font_pt=font_pt, bold=bold, color=color)


def _footnote(slide, text):
    """Small footnote at the bottom of the slide content area."""
    _txbox(slide, text,
           Inches(0.72), Inches(6.45), Inches(10.2), Inches(0.4),
           font_pt=8, color=GREY, italic=True)


def _rec_badge(slide, text="✔ Recommended", left=Inches(0.91), top=Inches(1.0)):
    """Small 'Recommended' badge."""
    _txbox(slide, text, left, top, Inches(2.2), Inches(0.28),
           font_pt=9, bold=True, color=BLUE)


def _options_table(slide, rows, col_headers, top=Inches(1.1),
                   col_widths=None):
    """
    Draw a simple borderless table-like layout using text boxes.
    rows: list of lists of strings
    col_headers: list of header strings
    """
    n_cols = len(col_headers)
    total_w = Inches(11.5)
    if col_widths is None:
        col_widths = [total_w / n_cols] * n_cols
    left_start = Inches(0.91)
    col_lefts = []
    x = left_start
    for cw in col_widths:
        col_lefts.append(x)
        x += cw

    row_h = Inches(0.32)
    header_h = Inches(0.35)

    # Header row
    for c, (hdr, cw, cl) in enumerate(zip(col_headers, col_widths, col_lefts)):
        _txbox(slide, hdr, cl, top, cw - Inches(0.05), header_h,
               font_pt=10, bold=True, color=BLUE)

    # Data rows
    for r, row in enumerate(rows):
        row_top = top + header_h + r * row_h
        for c, (cell, cw, cl) in enumerate(zip(row, col_widths, col_lefts)):
            color = BLACK
            bold = False
            if cell.startswith("✔"):
                color = BLUE
                bold = True
            _txbox(slide, cell, cl, row_top, cw - Inches(0.05), row_h - Inches(0.02),
                   font_pt=10, color=color, bold=bold)


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _add_slide(prs, "Title Slide")
    # Title
    for ph in sl.placeholders:
        try:
            if ph.placeholder_format.idx == 0:
                ph.text = "AI Blueprint"
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = "Source Sans Pro"
                        r.font.size = Pt(36)
                        r.font.bold = True
        except Exception:
            pass
    # Subtitle (Operations)
    _txbox(sl, "Operations",
           Inches(3.31), Inches(3.47), Inches(7.92), Inches(0.84),
           font_pt=28, bold=False, color=BLUE)
    # Team / date
    _txbox(sl, "AO  |  May 2026",
           Inches(9.0), Inches(5.75), Inches(3.5), Inches(0.5),
           font_pt=12, color=GREY)
    _txbox(sl, "Design & Approach Validation",
           Inches(3.31), Inches(4.1), Inches(7.92), Inches(0.5),
           font_pt=14, color=DARK_GREY)


def slide_playing_rules(prs, slide_num):
    sl = _add_slide(prs)
    _title(sl, "Playing Rules")
    _tag(sl, "For Information")
    body = (
        "For Information – Slides marked with this tag present a decision IRAS has taken. "
        "Reading is optional; no response is required.\n"
        "\n"
        "For Confirmation – Slides marked with this tag present an option or approach that IRAS "
        "is seeking external participants' confirmation on. "
        "Participants are expected to review and provide comments or validation.\n"
        "\n"
        "For Inputs – Slides marked with this tag are open questions. "
        "Participants are encouraged to share recommendations, suggestions, or solutions.\n"
        "\n"
        "Where multiple options are presented, IRAS' recommended option is marked ✔."
    )
    _body(sl, body, Inches(0.91), height=Inches(5.2), font_pt=12)


def slide_agenda(prs):
    sl = _add_slide(prs)
    _title(sl, "Topic Coverage / Scope Validation")
    _tag(sl, "For Confirmation")

    covered = (
        "Release Management\n"
        "  • Agent Fleet & Version Registry\n"
        "  • Deployment Lifecycle (Dev → Test → Pre-Prod → Prod)\n"
        "  • Testing Strategy (Blue-Green + Canary)\n"
        "\n"
        "Monitoring, Observability & Auditing\n"
        "  • EvalOps: Continuous Quality Evaluation\n"
        "  • Distributed Tracing (OpenTelemetry + Langfuse)\n"
        "  • Behavioural Auditing & Anomaly Detection\n"
        "\n"
        "Incident Management\n"
        "  • Human-in-the-Loop (HITL) Management\n"
        "  • Emergency Stop & Kill Switch Hierarchy"
    )
    covered2 = (
        "Platform Reliability & SLAs\n"
        "  • Circuit Breakers & Fallback Routing\n"
        "  • SLA Framework (Speed, Quality, Availability, Cost)\n"
        "\n"
        "FinOps: Token & Cost Management\n"
        "  • Per-session cost attribution & guardrails\n"
        "  • Cost comparison: Agent vs. Human Officer\n"
        "\n"
        "Operational Feedback & Lifecycle\n"
        "  • Structured feedback loop\n"
        "  • Model / application lifecycle management\n"
        "\n"
        "Operating Model\n"
        "  • Roles & responsibilities (RACI)\n"
        "  • Runbooks & escalation playbooks"
    )
    _txbox(sl, covered,
           Inches(0.91), Inches(0.75), Inches(5.6), Inches(5.8),
           font_pt=11, color=BLACK)
    _txbox(sl, covered2,
           Inches(6.8), Inches(0.75), Inches(5.6), Inches(5.8),
           font_pt=11, color=BLACK)

    _txbox(sl, "For Confirmation – Has IRAS covered all the Operations topics for AO that should be included in the Blueprint?",
           Inches(0.91), Inches(6.1), Inches(11.5), Inches(0.5),
           font_pt=11, bold=True, color=RED)


def slide_platform_arch(prs):
    sl = _add_slide(prs)
    _title(sl, "Platform Architecture Overview (AO Perspective)")
    _tag(sl, "For Information")

    arch = (
        "The AIBP platform processes taxpayer emails through four sequential layers. "
        "Operations touches every layer.\n"
        "\n"
        "Taxpayer Email\n"
        "    ↓\n"
        "SWEE  (Azure App Service / ACA, GCC 2.0)\n"
        "  • Receive inbound email  • Anonymise / PII-strip\n"
        "  • Classify to SOP via vector search  • Enqueue to Service Bus\n"
        "    ↓  Azure Service Bus Premium\n"
        "Agentic Orchestration – AO  (Azure Container Apps, LangGraph)\n"
        "  • One Container App per agent (one-agent-per-container architecture)\n"
        "  • Dequeue, execute agent workflow, call tools via AIGP\n"
        "    ↓  AIGP API (ForgeRock AuthZ + Entra ID AuthN)\n"
        "AI Governance Platform – AIGP\n"
        "  • Policy enforcement (OPA)  • HITL gating  • Audit trail\n"
        "    ↓  Kafka (AKS-hosted)\n"
        "Internal Microservices / DB"
    )
    _body(sl, arch, Inches(0.91), height=Inches(5.6), font_pt=11)


def slide_op_principles(prs):
    sl = _add_slide(prs)
    _title(sl, "Operational Principles")
    _tag(sl, "For Information")

    principles = (
        "1.  Observability-first\n"
        "    No agent may be deployed to production without emitting structured telemetry. "
        "If it cannot be observed, it cannot be operated.\n"
        "\n"
        "2.  Fail safe, not silent\n"
        "    A failed or uncertain agent action must escalate to a human. "
        "An agent that silently succeeds on the wrong action is worse than one that fails loudly.\n"
        "\n"
        "3.  Data minimisation in ops tooling\n"
        "    Operational tooling must handle only anonymised or metadata-level data. "
        "PII must not appear in any ops pipeline, log stream, or dashboard.\n"
        "\n"
        "4.  Policy as code\n"
        "    Operational guardrails (rate limits, kill switches, risk thresholds) are "
        "version-controlled code. Manual configuration drift is a reliability and audit risk.\n"
        "\n"
        "5.  Separation of concerns\n"
        "    AO team owns agent logic; AIGP team owns governance enforcement; Ops team owns "
        "production health. No single team has unilateral access to all layers in production."
    )
    _body(sl, principles, Inches(0.91), height=Inches(5.8), font_pt=11)


# ── Release Management ─────────────────────────────────────────────────────────

def slide_agent_registry(prs):
    sl = _add_slide(prs)
    _title(sl, "Agent Fleet & Version Registry")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Release Management", Inches(0.6))

    _txbox(sl, "Options Comparison",
           Inches(0.91), Inches(0.75), Inches(11.5), Inches(0.35),
           font_pt=13, bold=True, color=BLACK)

    headers = ["Option", "Approach", "Pros", "Cons"]
    col_widths = [Inches(1.3), Inches(3.8), Inches(3.2), Inches(3.2)]
    rows = [
        ["✔ Option 1\n(Recommended)",
         "SemVer (MAJOR.MINOR.PATCH) + ACR Signed Images (Notation/CNCF) "
         "+ PostgreSQL Registry DB\n"
         "• agent-manifest.json per agent repo\n"
         "• Append-only ledger (no UPDATE/DELETE)\n"
         "• Image signed via Azure Key Vault",
         "• Cryptographic guarantee: what was tested in CI runs in prod\n"
         "• Permanent, auditable version history\n"
         "• Targeted rollback possible (PATCH without losing MINOR)\n"
         "• ISO 27001 CI-04 compliant",
         "• PostgreSQL registry adds infra dependency\n"
         "• PKI cert rotation must be planned for"],
        ["Option 2",
         "Git SHA tagging + flat JSON deployment file",
         "• Zero additional infrastructure",
         "• Opaque to ops tooling; no eval score association\n"
         "• Cannot satisfy ISO 27001 audit requirements"],
        ["Option 3",
         "Monolithic release (all agents versioned together)",
         "• Simpler release coordination",
         "• Any change forces full redeployment\n"
         "• Regression in one agent rolls back all agents"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.1), col_widths=col_widths)

    _txbox(sl, "IRAS Recommendation: Option 1 — append-only registry + signed images provides the operational forensics "
           "and ISO 27001 change management audit trail required.",
           Inches(0.91), Inches(5.9), Inches(11.5), Inches(0.5),
           font_pt=10, bold=True, color=BLUE)


def slide_deployment_lifecycle(prs):
    sl = _add_slide(prs)
    _title(sl, "Deployment Lifecycle: Dev → Test → Pre-Prod → Prod")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Release Management", Inches(0.6))

    _txbox(sl, "Options Comparison",
           Inches(0.91), Inches(0.75), Inches(11.5), Inches(0.35),
           font_pt=13, bold=True, color=BLACK)

    headers = ["Option", "Approach", "Pros", "Cons"]
    col_widths = [Inches(1.3), Inches(3.8), Inches(3.2), Inches(3.2)]
    rows = [
        ["✔ Option 1\n(Recommended)",
         "Azure DevOps Multi-Stage Pipeline with Manual Approval Gates\n"
         "• Stage 1: Build & Sign (PR merge → main)\n"
         "• Stage 2: Test env + Synthetic Email Eval Suite (100 labelled emails)\n"
         "• Stage 3: Pre-Prod — manual gate (AO team lead approval)\n"
         "• Stage 4: Prod Canary — double gate (AO lead + Ops)",
         "• All gates recorded in Azure DevOps + registry DB\n"
         "• ISO 27001 change audit trail\n"
         "• Environment isolation (separate ACA env, SB namespace, AOAI quota pool)\n"
         "• Double approval prevents unilateral prod deploy",
         "• Patch can take 24–72 h from merge to full prod (approval turnaround)\n"
         "• Synthetic eval dataset requires ongoing curation"],
        ["Option 2",
         "GitHub Actions + Environment Protection Rules",
         "• Better developer experience; native PR integration",
         "• Self-hosted runners required in GCC 2.0\n"
         "• Introduces second CI/CD platform to govern"],
        ["Option 3",
         "Manual deployment with runbook",
         "• No CI/CD infra to maintain",
         "• No automated eval gate\n"
         "• Fails ISO 27001 change management\n"
         "• Image signing skipped under time pressure"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.1), col_widths=col_widths)

    _txbox(sl, "IRAS Recommendation: Option 1 — Azure DevOps is within existing GCC 2.0 governance boundary; "
           "the pipeline encodes eval thresholds as automated enforcement, not a policy document.",
           Inches(0.91), Inches(5.9), Inches(11.5), Inches(0.5),
           font_pt=10, bold=True, color=BLUE)


def slide_testing_strategy(prs):
    sl = _add_slide(prs)
    _title(sl, "Testing Strategy: Blue-Green + Canary Deployment")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Release Management", Inches(0.6))

    body = (
        "Options Comparison\n"
        "\n"
        "✔ Option 1 (Recommended) — Blue-Green with Canary (ACA Traffic Weight Splitting)\n"
        "  • New agent version deployed as a new ACA revision; initially receives 10% of live traffic\n"
        "  • Canary window: 24–48 hours; promotion criteria checked before full rollout\n"
        "  • Canary promotion criteria: error rate ≤ blue, P95 latency ≤ blue +10%, "
        "eval composite score ≥ blue –2%, no behavioral anomaly alerts\n"
        "  • Rollback: set canary traffic weight to 0%; no redeployment needed\n"
        "  Pros: Real taxpayer email exposure at low volume; non-deterministic LLM behaviour "
        "surfaces only in live data; instant rollback\n"
        "  Cons: Requires ACA traffic splitting infrastructure; two versions process emails "
        "simultaneously during canary window\n"
        "\n"
        "Option 2 — Blue-Green Switch (All-or-Nothing Cutover)\n"
        "  Cons: Too risky for LLMs; subtle hallucinations may not surface in staging; "
        "100% of emails immediately on new version\n"
        "\n"
        "Option 3 — Shadow Mode (Parallel Run without Serving)\n"
        "  Pros: No risk to taxpayer emails during testing\n"
        "  Cons: Shadow traffic does not surface all production failure modes; "
        "does not validate tool calls that require real data responses\n"
        "\n"
        "IRAS Recommendation: Option 1. LLMs are non-deterministic; "
        "the 10% canary at low volume is the safest path to real-world validation."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.8), font_pt=10.5)


# ── Observability ──────────────────────────────────────────────────────────────

def slide_evalops(prs):
    sl = _add_slide(prs)
    _title(sl, "EvalOps: Continuous Quality Evaluation")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Monitoring & Observability", Inches(0.6))

    body = (
        "Options Comparison\n"
        "\n"
        "✔ Option 1 (Recommended) — Azure AI Evaluation SDK + LLM-as-a-Judge (GPT-4o) + Langfuse (self-hosted)\n"
        "  Evaluation metrics per agent run:\n"
        "  • Faithfulness (groundedness): % of response claims supported by retrieved context\n"
        "  • Answer relevance: alignment between response and the taxpayer's query\n"
        "  • Task completion: did the agent call the correct tools in the correct order?\n"
        "  • Hallucination rate: unsupported factual claims flagged by GPT-4o judge\n"
        "  • Tool-call accuracy (deterministic): pytest assertions on tool call sequence\n"
        "  Evaluation runs in CI (per PR) and in production (sampled real traces via Langfuse)\n"
        "  CI gate: block promotion if composite score drops >5% vs current active version\n"
        "  Pros: Scalable, automated, catches semantic regressions not visible in unit tests\n"
        "  Cons: LLM-as-judge has slight verbosity bias; judge model cost (~USD 10–30/month)\n"
        "\n"
        "Option 2 — Deterministic Unit Testing Only\n"
        "  Pros: Cheap and fast   Cons: Cannot catch reasoning quality regressions or hallucinations\n"
        "\n"
        "Option 3 — Human Officer Spot-Audit Only\n"
        "  Cons: Reactive; errors are found after taxpayer impact; not scalable\n"
        "\n"
        "IRAS Recommendation: Option 1, supplemented by Option 2 for tool-call structural checks."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


def slide_observability(prs):
    sl = _add_slide(prs)
    _title(sl, "Observability & Distributed Tracing")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Monitoring & Observability", Inches(0.6))

    body = (
        "Options Comparison\n"
        "\n"
        "✔ Option 1 (Recommended) — OpenTelemetry (OTel) + OpenInference Semantic Conventions "
        "→ Azure Monitor Application Insights + Langfuse (self-hosted)\n"
        "\n"
        "  What is instrumented per email processing event:\n"
        "  • Root span: messageId, agentName, agentSemver, sopId, resolutionType\n"
        "  • Child spans per LLM call: model, prompt token count, completion token count, latency\n"
        "  • Child spans per tool call: tool name, input params (anonymised), output summary, latency\n"
        "  • Custom metrics: ao.email.latency_ms, ao.email.tokens.prompt/completion, "
        "ao.email.tool_calls.count, ao.email.hitl_escalated\n"
        "\n"
        "  Azure Monitor Workbooks surface:\n"
        "  • AO Operations Dashboard: P50/P95 latency, token consumption trends, HITL rate\n"
        "  • Canary vs. Blue revision side-by-side comparison during deployment\n"
        "\n"
        "  PII control: all spans are anonymised before export; taxpayer email body "
        "and personal data are never in telemetry. Trace IDs link back to the encrypted blob.\n"
        "\n"
        "Option 2 — LangSmith (LangChain hosted)\n"
        "  Cons: External SaaS; data residency in GCC 2.0 not guaranteed\n"
        "\n"
        "Option 3 — Standard application logs only\n"
        "  Cons: Cannot reconstruct agent reasoning chain; hallucinations are undebuggable post-mortem\n"
        "\n"
        "IRAS Recommendation: Option 1. Self-hosted Langfuse keeps trace data within GCC 2.0 boundary."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


def slide_behavioural_audit(prs):
    sl = _add_slide(prs)
    _title(sl, "Behavioural Auditing, Traceability & Anomaly Detection")
    _tag(sl, "For Information")
    _section_label(sl, "Monitoring & Observability", Inches(0.6))

    body = (
        "Behavioural Audit (Append-Only Cosmos DB Ledger)\n"
        "  Every agent decision step is written to an append-only Cosmos DB audit ledger:\n"
        "  • Fields: messageId, agentName, agentSemver, stepType, inputSummary (anonymised), "
        "outputSummary (anonymised), toolName, riskScore, aigpDecision, timestamp\n"
        "  • Ledger is write-once; no UPDATE or DELETE is permitted (PDPA / IM8 compliant)\n"
        "  • Anonymised telemetry stream emitted to Event Hubs for SOC consumption\n"
        "  • Audit ledger retained for 7 years (IM8 / WOG record retention requirement)\n"
        "\n"
        "Behavioural Anomaly Detection (Embedding-Based Baseline)\n"
        "  • During CI evaluation, a 'normal behavioural baseline' is generated: "
        "each agent's tool-call sequences across synthetic emails are embedded using "
        "text-embedding-3-small and stored as a reference distribution\n"
        "  • In production, each completed trace is embedded and its cosine distance "
        "from the baseline distribution is computed\n"
        "  • Alert thresholds:\n"
        "      – >20% of traces deviate significantly → P2 alert, AO team notified\n"
        "      – >50% deviation rate → automated Emergency Stop triggered\n"
        "\n"
        "Note on CoT visibility: Azure OpenAI gpt-4o does not expose internal "
        "chain-of-thought reasoning tokens in the API response. "
        "The audit ledger captures the observable inputs and outputs at each LangGraph node step, "
        "providing a reconstructable decision trail without requiring internal CoT access."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.8), font_pt=10.5)


# ── Incident Management ────────────────────────────────────────────────────────

def slide_hitl(prs):
    sl = _add_slide(prs)
    _title(sl, "Human-in-the-Loop (HITL) Management")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Incident Management", Inches(0.6))

    body = (
        "AIGP Risk-Score-Based Escalation (Risk score 0.00 – 1.00)\n"
    )
    _txbox(sl, body, Inches(0.91), Inches(0.75), Inches(11.5), Inches(0.35),
           font_pt=13, bold=True)

    headers = ["Risk Score", "AIGP Decision", "Example Actions"]
    col_widths = [Inches(2.0), Inches(3.0), Inches(6.5)]
    rows = [
        ["0.00–0.40",   "AUTO-PERMIT",         "aigp.get_taxpayer_record (read-only, low risk)"],
        ["0.41–0.69",   "PERMIT WITH AUDIT",   "aigp.submit_refund_request ≤ SGD 10,000"],
        ["0.70–0.89",   "HITL REQUIRED",       "aigp.close_case, aigp.submit_refund_request > SGD 10,000"],
        ["0.90–1.00",   "AUTO-BLOCK",          "aigp.update_bank_details (irreversible; fraud vector)"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.15), col_widths=col_widths)

    _txbox(sl, "HITL Review App Options",
           Inches(0.91), Inches(2.85), Inches(11.5), Inches(0.32),
           font_pt=13, bold=True)

    body2 = (
        "✔ Option 1A (Recommended) — Custom React webapp on Azure Static Web Apps + Azure Functions API\n"
        "  Task inbox sorted by SLA deadline | Risk factor badges | Deep link to case management system "
        "| Approve/Reject with mandatory notes (min 20 chars) | SLA countdown timer\n"
        "  Pros: Purpose-built UX; fully within GCC 2.0; officer can complete review in <2 min\n"
        "  Cons: 2–3 weeks build effort for MVP\n"
        "\n"
        "Option 1B — Power Apps Canvas App + Power Automate\n"
        "  Pros: Faster initial build; Teams push notifications\n"
        "  Cons: Power Platform GCC 2.0 availability must be confirmed; SharePoint task store "
        "does not natively integrate with Service Bus\n"
        "\n"
        "HITL SLA Tiers:\n"
        "  Tier 1 (Standard, risk 0.70–0.84): 4 business hours → auto-escalate to supervisor\n"
        "  Tier 2 (Urgent, risk 0.85–0.89 or disputed taxpayer): 1 business hour → escalate to branch manager\n"
        "  Capacity at <1,000 emails/day + 10–15% HITL rate: ~100–150 HITL tasks/day "
        "(~8–25 officer-hours/day)"
    )
    _body(sl, body2, Inches(3.2), height=Inches(3.3), font_pt=10)


def slide_emergency_stop(prs):
    sl = _add_slide(prs)
    _title(sl, "Incident Response & Emergency Stop (Kill Switch Hierarchy)")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Incident Management", Inches(0.6))

    body = (
        "Kill Switch Hierarchy (3 levels, all via Azure App Configuration feature flags)\n"
        "\n"
        "  Level 1 — Agent-Version Stop: agents:{agentName}:{semver}:enabled = false\n"
        "            Scope: one specific version; other versions continue running\n"
        "\n"
        "  Level 2 — Agent-All-Versions Stop: agents:{agentName}:enabled = false\n"
        "            Scope: all versions of one agent; other agents unaffected\n"
        "\n"
        "  Level 3 — Platform Stop: agents:platform:enabled = false\n"
        "            Scope: all agents; Service Bus topic disabled to halt new enqueueing\n"
        "\n"
        "Trigger paths:\n"
        "  • Automated: Behavioural anomaly >50% deviation rate triggers Level 1/2 automatically\n"
        "  • Manual: Ops Lead (Level 1/2) or Ops Lead + AIGP Team Lead (Level 3)\n"
        "  • Response time: Flag propagation to all ACA containers within 90 seconds\n"
        "  • In-flight messages: Lock expiry returns messages to queue for retry/DLQ\n"
        "\n"
        "Platform Resume (after Emergency Stop)\n"
        "  Requires double approval: Ops Lead + AO Team Lead sign-off in Azure DevOps\n"
        "  AO team must confirm root cause and fix deployed before resume is approved\n"
        "  All stop/resume events written to append-only Emergency Stop Cosmos DB container\n"
        "\n"
        "Fallback path during any outage → emails routed to Human Officer Queue "
        "(unified escalation queue for all categories the platform cannot process automatically)\n"
        "\n"
        "This constitutes the Business Continuity Plan (BCP) for the AO layer (IM8 requirement): "
        "manual processing capacity must be verified sufficient for full email volume."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


# ── Platform Reliability ───────────────────────────────────────────────────────

def slide_reliability(prs):
    sl = _add_slide(prs)
    _title(sl, "Platform Reliability: Circuit Breakers & Fallback Routing")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Platform Reliability & SLAs", Inches(0.6))

    _txbox(sl, "Options Comparison",
           Inches(0.91), Inches(0.75), Inches(11.5), Inches(0.35),
           font_pt=13, bold=True)

    headers = ["Option", "Approach", "Pros", "Cons"]
    col_widths = [Inches(1.4), Inches(4.2), Inches(2.85), Inches(3.05)]
    rows = [
        ["✔ Option 1\n(Recommended)",
         "Tenacity Circuit Breakers (Python) + Service Bus DLQ → Human Officer Queue fallback "
         "+ Azure OpenAI Provisioned Throughput Units (PTU)\n"
         "• Circuit states: Closed / Open / Half-Open per dependency\n"
         "• Fallback: DLQ → human-officer-queue → Tax officer notified via Teams\n"
         "• 25 PTU (gpt-4o) eliminates PAYG rate-limit throttling as a failure mode",
         "• Tenacity is battle-tested; no extra infra\n"
         "• DLQ → human queue ensures 100% email durability\n"
         "• PTU removes most common failure mode",
         "• PTU is upfront monthly commitment\n"
         "• Circuit state is per-container (not shared across ACA replicas)"],
        ["Option 2",
         "PAYG Azure OpenAI + Static Retry Logic (no circuit breaker)",
         "• No PTU commitment; simpler implementation",
         "• Retry storms during throttle events compound outage\n"
         "• PAYG latency variance causes SLA breach risk"],
        ["Option 3",
         "Active-Passive AO Cluster (secondary ACA environment in second region)",
         "• Maximum resiliency against full region failure",
         "• Doubles infra cost for very low-probability scenario\n"
         "• Not recommended at <1,000 emails/day scale"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.1), col_widths=col_widths)

    _txbox(sl, "IRAS Recommendation: Option 1. At <1,000 emails/day, circuit breakers + PTU + DLQ fallback "
           "cover the two highest-impact failure modes at appropriate cost.",
           Inches(0.91), Inches(5.9), Inches(11.5), Inches(0.5),
           font_pt=10, bold=True, color=BLUE)


def slide_sla(prs):
    sl = _add_slide(prs)
    _title(sl, "SLA Framework — Speed, Quality, Availability, Cost")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Platform Reliability & SLAs", Inches(0.6))

    _txbox(sl, "Proposed SLA targets (greenfield baseline — to be validated in first 90 days of production)",
           Inches(0.91), Inches(0.72), Inches(11.5), Inches(0.32),
           font_pt=11, bold=True, color=BLACK)

    headers = ["Dimension", "KPI", "Proposed Target", "Measurement"]
    col_widths = [Inches(1.5), Inches(3.8), Inches(2.6), Inches(3.6)]
    rows = [
        ["Speed",        "Time to Resolution (Automated)",         "P95 ≤ 5 min (biz hours)",    "ao.email.latency_ms (OTel)"],
        ["Speed",        "Time to Resolution (HITL Tier 1)",       "P95 ≤ 4 biz hours",          "HITL task lifecycle record"],
        ["Speed",        "Time to Resolution (HITL Tier 2/Urgent)", "P95 ≤ 1 biz hour",          "HITL task lifecycle record"],
        ["Quality",      "Automated Resolution Accuracy",           "≥ 95%",                      "hitl.correctionRate (inverse)"],
        ["Quality",      "Hallucination Rate",                      "≤ 3%",                       "EvalOps LLM-as-Judge (faithfulness)"],
        ["Quality",      "Misroute Rate",                           "≤ 5%",                       "ao.email.rerouteRate"],
        ["Availability", "Platform Availability (biz hours)",       "≥ 99.5%",                    "Azure Monitor heartbeat"],
        ["Availability", "DLQ Rate",                                "≤ 1%",                       "servicebus.dlq.messageCount"],
        ["Cost",         "Cost per Automated Resolution",           "≤ SGD 0.20",                 "FinOps attribution (Step 5)"],
        ["Cost",         "Automation Rate",                         "≥ 80% (90-day); ≥ 90% (1yr)","resolutionType metric"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.08), col_widths=col_widths)

    body2 = (
        "SLA Dashboard: Azure Monitor Workbook (real-time RAG status per KPI) | "
        "Monthly SLA Report: auto-generated via Azure Logic Apps\n"
        "Option 2 (Manual spreadsheet): No real-time visibility; 4–8 h analyst effort/month — not recommended\n"
        "Option 3 (Power BI): Requires GCC 2.0 Power BI availability confirmation; 15-min refresh minimum vs. real-time Workbooks"
    )
    _txbox(sl, body2, Inches(0.91), Inches(5.85), Inches(11.5), Inches(0.55),
           font_pt=9.5, color=DARK_GREY)


# ── FinOps ─────────────────────────────────────────────────────────────────────

def slide_finops(prs):
    sl = _add_slide(prs)
    _title(sl, "FinOps: Token & Cost Management")
    _tag(sl, "For Confirmation")
    _section_label(sl, "FinOps", Inches(0.6))

    _txbox(sl, "Options Comparison",
           Inches(0.91), Inches(0.75), Inches(11.5), Inches(0.35),
           font_pt=13, bold=True)

    headers = ["Option", "Approach", "Pros", "Cons"]
    col_widths = [Inches(1.4), Inches(4.0), Inches(3.1), Inches(3.0)]
    rows = [
        ["✔ Option 1\n(Recommended)",
         "Per-session cost attribution via OTel custom metrics + Azure Cost Management + automated guardrails\n"
         "• Every email: PromptTokens, CompletionTokens, TotalLLMCostUSD, ComputeCostUSD stored in "
         "Log Analytics (AIBPEmailCosts table)\n"
         "• Per-session token budget enforced at agent level: flag set in Azure App Configuration; "
         "agent stops calling LLM and routes to Human Officer Queue if exceeded\n"
         "• Daily + monthly budget alerts at 90% / 100% of budget\n"
         "• Estimated cost: ~USD 1,410/month for <1,000 emails/day (25 PTU gpt-4o flat rate)",
         "• Per-email attribution identifies which agent / SOP is expensive\n"
         "• Per-session guardrail prevents runaway agent loops\n"
         "• No email is silently costing 50× average",
         "• FinOps aggregator ACA Job adds maintenance overhead\n"
         "• PTU flat rate may be higher than PAYG at very low utilisation"],
        ["Option 2",
         "Azure Cost Management aggregate monitoring only (no per-session attribution)",
         "• Zero implementation effort",
         "• Cannot identify which agent or email is expensive\n"
         "• No per-session guardrail; stuck agent burns budget invisibly"],
        ["Option 3",
         "OpenCost on AKS (Kubernetes-level cost attribution)",
         "• Provides AKS compute attribution",
         "• Does not cover Azure OpenAI token cost (the dominant cost driver)\n"
         "• Adds infra complexity for a secondary cost category"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.1), col_widths=col_widths)

    _txbox(sl, "IRAS Recommendation: Option 1. Token cost varies dramatically per email; "
           "per-session attribution + guardrail is the single most important cost control in an agentic platform.",
           Inches(0.91), Inches(5.85), Inches(11.5), Inches(0.5),
           font_pt=10, bold=True, color=BLUE)


def slide_cost_comparison(prs):
    sl = _add_slide(prs)
    _title(sl, "Cost Comparison: Agent-Resolved vs. Human-Officer-Resolved")
    _tag(sl, "For Information")
    _section_label(sl, "FinOps", Inches(0.6))

    body = (
        "Human Officer Cost Per Email (illustrative — to be validated with HR & Finance)\n"
        "  • Officer annual salary: SGD A (to be provided by HR)\n"
        "  • Employer overhead (CPF, benefits, office): ~30–40% additional\n"
        "  • Illustrative: SGD 80,000 salary + 35% overhead = SGD 108,000/year\n"
        "  • Emails processed/officer/year: ~17,500 (70/day × 250 days) — routine cases\n"
        "  • Illustrative human cost per email: SGD 108,000 ÷ 17,500 ≈ SGD 6.17/email\n"
        "\n"
        "Agent Cost Per Email (from FinOps attribution, Step 5.1)\n"
        "  • Azure OpenAI token cost: ~USD 0.009 (~SGD 0.012) per email at 3,000 avg tokens\n"
        "  • ACA compute cost: ~SGD 0.003 per email at P95 5-min processing\n"
        "  • Overhead (Langfuse, monitoring, registry DB): ~SGD 0.005 per email\n"
        "  • Total agent cost per email: ~SGD 0.02–0.10 (to be validated in first 30 days)\n"
        "  • Target: ≤ SGD 0.20/resolution\n"
        "\n"
        "Break-Even / ROI Analysis\n"
        "  At 80% automation rate (1,000 emails/day → 800 automated, 200 HITL/human):\n"
        "  • Human cost equivalent: 800 × SGD 6.17 = SGD 4,936/day\n"
        "  • Agent cost: 800 × SGD 0.10 = SGD 80/day\n"
        "  • Net daily saving (automated portion): ~SGD 4,856/day\n"
        "  • Platform monthly operating cost: ~SGD 1,890 (Azure OpenAI) + other infra\n"
        "\n"
        "Note: These are indicative estimates. Ops team to establish actuals within first 30 days of production "
        "and publish a validated cost comparison to management."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


# ── Feedback & Lifecycle ───────────────────────────────────────────────────────

def slide_feedback_loop(prs):
    sl = _add_slide(prs)
    _title(sl, "Operational Feedback Loop & Self-Reflection")
    _tag(sl, "For Confirmation")
    _section_label(sl, "Operational Feedback & Lifecycle", Inches(0.6))

    body = (
        "Options Comparison\n"
        "\n"
        "✔ Option 1 (Recommended) — Event-Driven Feedback Pipeline → Langfuse Evaluation Dataset\n"
        "\n"
        "  Signal sources (4):\n"
        "  1. HITL officer decisions (approve/reject) — high rejection rate for an action type "
        "signals the agent is over-proposing it\n"
        "  2. Officer audit corrections (SWEE triage) — direct accuracy signal for the triage model\n"
        "  3. Human Officer Queue cases (fallback) — cases the platform failed to handle; "
        "officer resolution becomes a labelled training example\n"
        "  4. AO agent re-route events — agent explicitly rejected SWEE's SOP; "
        "triage accuracy signal\n"
        "\n"
        "  Pipeline: Azure Service Bus topic (ops-feedback-events) → "
        "feedback-processor ACA Job (runs every 15 min) → "
        "annotates Langfuse traces with structured scores (0.0 = failure, 1.0 = correct)\n"
        "  Failure traces added to Langfuse evaluation datasets:\n"
        "    • swee-triage-failures (for SWEE team — weekly review)\n"
        "    • ao-hitl-rejections (for AO team — biweekly review)\n"
        "    • ao-capability-gaps (for AO team — monthly review)\n"
        "  Quality control: 30-day quarantine gate before failures enter CI evaluation baseline\n"
        "  Ops sends weekly Feedback Summary Report to owning teams (every Monday 09:00)\n"
        "\n"
        "Option 2 — Weekly manual feedback collation by Ops\n"
        "  Cons: Up to 7-day lag; 70–100 emails can be affected before team is notified; "
        "does not produce Langfuse dataset annotations\n"
        "\n"
        "IRAS Recommendation: Option 1. Langfuse dataset annotations create a compounding quality asset "
        "that makes the CI evaluation gate progressively more sensitive to real-world failure patterns."
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


def slide_model_lifecycle(prs):
    sl = _add_slide(prs)
    _title(sl, "Model / Application Lifecycle Management")
    _tag(sl, "For Information")
    _section_label(sl, "Operational Feedback & Lifecycle", Inches(0.6))

    body = (
        "Model Upgrade Strategy\n"
        "  • Azure OpenAI model upgrades (e.g., gpt-4o → gpt-4o-next) are treated as a "
        "MAJOR agent version bump — full CI pipeline + canary deployment required\n"
        "  • Evaluation thresholds in agent-manifest.json must be re-validated against the new model "
        "before the first production canary\n"
        "  • PTU allocation may need re-sizing for new model; review pricing and TPM characteristics\n"
        "\n"
        "Agent Deprecation Process\n"
        "  • When an SOP category is retired or merged, the corresponding agent version is marked "
        "deployment_status = 'deprecated' in the registry DB\n"
        "  • Service Bus subscription for that SOP category is disabled (not deleted — "
        "for audit trail retention)\n"
        "  • Deprecation is announced to the SWEE team to remove the SOP from the triage routing\n"
        "\n"
        "SOP Corpus Updates\n"
        "  • Tax policy changes (new SOPs, amended SOPs) require:\n"
        "    1. SWEE team to update the SOP vector index in Azure AI Search\n"
        "    2. AO team to review if the affected SOP's agent requires a prompt update (PATCH version)\n"
        "    3. New agent version goes through the full CI/CD pipeline\n"
        "\n"
        "Automated Post-Mortem Pipeline\n"
        "  • Triggered by: Emergency Stop, behavioural anomaly >50%, HITL SLA breach >10%, "
        "accuracy drop >5%, DLQ spike >3%\n"
        "  • postmortem-assembler ACA Job: reconstructs incident timeline from Langfuse + Azure Monitor; "
        "retrieves 20 most anomalous traces; GPT-4o summarises the pattern\n"
        "  • Output: auto-drafted post-mortem report in Azure Blob Storage, distributed via Teams\n"
        "  • Human analyst completes the causal analysis ('why') and remediation plan\n"
        "  • Post-mortem must be completed within 5 business days of the incident"
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


# ── Operating Model ────────────────────────────────────────────────────────────

def slide_raci(prs):
    sl = _add_slide(prs)
    _title(sl, "Operating Model: Roles & Responsibilities (RACI)")
    _tag(sl, "For Information")
    _section_label(sl, "Operating Model", Inches(0.6))

    _txbox(sl, "Team Ownership Summary",
           Inches(0.91), Inches(0.72), Inches(11.5), Inches(0.32),
           font_pt=12, bold=True)

    headers = ["Team", "Owns", "Day-to-Day Responsibilities"]
    col_widths = [Inches(2.0), Inches(3.8), Inches(5.7)]
    rows = [
        ["Ops Team",          "Platform health, SLA reporting, FinOps, HITL queue, Emergency Stop, feedback pipeline, post-mortems",
         "Monitor dashboards, respond to alerts, manage HITL SLA, generate reports, maintain compliance evidence"],
        ["AO Team\n(App & Data)", "Agent source code, manifest, CI/CD pipeline, agent eval logic",
         "Develop/deploy agent updates, respond to CI eval failures, review feedback datasets, release patches"],
        ["SWEE Team\n(App & Data)", "SWEE code, SOP vector index, email ingestion, triage model",
         "Monitor triage accuracy, respond to embedding drift, review swee-triage-failures dataset weekly"],
        ["AIGP Team",         "AIGP API, OPA policy, risk scoring config, HITL thresholds, ForgeRock/Entra integration",
         "Review OPA policy effectiveness, manage HITL threshold calibration, respond to AIGP-triggered stops"],
        ["Tax Officers\n(Business)", "HITL review decisions, SWEE audit corrections, Human Officer Queue resolutions",
         "Process HITL tasks within SLA tiers, review SWEE audit sample, resolve manual queue emails"],
        ["SOC (Internal\n+ GovTech)", "Security monitoring of AIBP infrastructure",
         "Monitor MDC alerts, review anonymised Event Hubs telemetry stream; no access to Langfuse or audit ledger"],
    ]
    _options_table(sl, rows, headers, top=Inches(1.08), col_widths=col_widths)

    _footnote(sl, "Full RACI matrix available in the Operations Blueprint (aibp-ops-step7.md)")


def slide_ops_model(prs):
    sl = _add_slide(prs)
    _title(sl, "Operating Model: Staffing, Access Control & Continuity")
    _tag(sl, "For Information")
    _section_label(sl, "Operating Model", Inches(0.6))

    body = (
        "Ops Team Structure (at <1,000 emails/day steady state — no 24/7 NOC required)\n"
        "  • 1 × Ops Lead: owns runbooks, chairs weekly platform health review, signs off Emergency Stop resumes\n"
        "  • 2 × Ops Engineers: monitor dashboards, respond to alerts, manage HITL SLA, generate reports\n"
        "  • On-call rotation: 1 Ops Engineer/week; after-hours scope = P1 alerts only "
        "(platform availability, Emergency Stop). P2/P3 queued for business-hours response.\n"
        "\n"
        "Access Control Principles\n"
        "  • Managed Identities only — no service account passwords or API keys in application code\n"
        "  • Least privilege RBAC — each team has minimum Azure role assignments needed for their function\n"
        "  • No shared accounts — every human access tied to named individual's Entra ID UPN\n"
        "  • Production access is pipeline-mediated — no direct ACA/Service Bus write access for humans "
        "(outside Emergency Stop / runbook path)\n"
        "  • Quarterly RBAC access reviews — Logic Apps auto-exports role assignments; Ops Lead recertifies\n"
        "\n"
        "Service Ownership & Continuity\n"
        "  • Each agent has a named service owner in the AO team (designated in agent-manifest.json)\n"
        "  • Platform BCP: DLQ → Human Officer Queue fallback is the documented BCP for the AO layer (IM8)\n"
        "  • BCP verification: quarterly drills to confirm manual processing capacity is sufficient\n"
        "  • Transition-to-run: formal knowledge transfer from the project team to steady-state Ops "
        "covering all runbooks, dashboards, escalation contacts, and compliance evidence procedures\n"
        "\n"
        "HITL Capacity Planning\n"
        "  At <1,000 emails/day + 10–15% HITL rate: ~100–150 HITL tasks/day (~8–25 officer-hours/day)\n"
        "  Alert: hitl.officer.taskLoad > 15 tasks/officer/day → P3 alert to Ops for staffing review"
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


def slide_runbooks(prs):
    sl = _add_slide(prs)
    _title(sl, "Runbooks & Escalation Playbooks")
    _tag(sl, "For Information")
    _section_label(sl, "Operating Model", Inches(0.6))

    body = (
        "Runbooks are step-by-step procedures for common and critical operational tasks. "
        "All runbooks are stored in the Ops team wiki (Azure DevOps Wiki). "
        "They are reviewed after every incident (within 5 business days) and quarterly.\n"
        "\n"
        "RB-01: New Agent Version Production Deployment\n"
        "  Trigger: AO team requests promotion of pre-prod-tested version to production canary\n"
        "  Key steps: Verify registry DB eval_pass=true → Trigger Stage 4 pipeline → "
        "Verify 10% ACA traffic split → Review canary metrics at T+24h and T+48h → "
        "Double-approve promotion (Ops + AO lead) → Verify registry DB updated\n"
        "\n"
        "RB-02: Agent Version Rollback\n"
        "  Trigger: Canary metrics fail; P2 alert during canary; behavioural anomaly\n"
        "  Key steps: Identify rollback target version from registry DB → "
        "Trigger rollback pipeline (set canary to 0%) → Spot-check 5 Langfuse traces → "
        "Notify AO team with anomalous trace IDs → Create post-mortem draft\n"
        "\n"
        "RB-03: Emergency Stop Activation (Manual)\n"
        "  Trigger: Observed anomalous behaviour; security concern; management direction\n"
        "  Key steps: Determine stop scope (agent-version / agent / platform) → "
        "Set App Configuration flag to false → Verify within 90s on dashboard → "
        "Handle in-flight messages → Notify all teams via Teams → Record in Cosmos DB stop ledger\n"
        "\n"
        "RB-04: DLQ Recovery\n"
        "  Trigger: DLQ rate >1% sustained; post-Emergency Stop backlog recovery\n"
        "  Key steps: Query DLQ for failure categories → Route 'AO_UNAVAILABLE' failures to "
        "Human Officer Queue → Replay resolvable failures after fix is confirmed deployed\n"
        "\n"
        "Escalation path: Ops Engineer → Ops Lead → AO/AIGP Team Lead → Platform Owner/CIO"
    )
    _body(sl, body, Inches(0.75), height=Inches(5.9), font_pt=10.5)


# ── Open Questions ─────────────────────────────────────────────────────────────

def slide_open_questions(prs):
    sl = _add_slide(prs)
    _title(sl, "Open Questions & Inputs Needed from External Participants")
    _tag(sl, "For Inputs")

    questions = (
        "Observability & Logs\n"
        "  Q1. What is the recommended minimum observability bar for any agentic workflow deployed in a "
        "government context (logs, traces, decision snapshots)? Who should have access for optimisation vs. audit?\n"
        "\n"
        "  Q2. If we redact all PII from traces and logs, we lose visibility needed for debugging. "
        "How do you strike the right balance between PDPA data minimisation and operational deep-visibility?\n"
        "\n"
        "  Q3. For OpenAI models where internal chain-of-thought tokens are not returned in the API response, "
        "what practices have you used to reconstruct 'why the agent made a specific decision' for an audit trail?\n"
        "\n"
        "  Q4. Frameworks for observability: Langfuse (self-hosted), LangSmith, Arize Phoenix, "
        "Pydantic Logfire — what has worked well in production? Any specific GCC 2.0 / data residency constraints?\n"
        "\n"
        "Release & Operations\n"
        "  Q5. For agentic workflows requiring frequent iterative LLM calls, how have you managed "
        "rate limits and token quotas in Azure enterprise cloud environments? "
        "Have you implemented semantic caching or request queuing at the infrastructure layer?\n"
        "\n"
        "  Q6. What is your recommended strategy for transitioning from a single-prompt RAG application "
        "to a multi-agent system in production — shadow mode concurrently, or migrate use cases one by one?\n"
        "\n"
        "  Q7. For non-deterministic failures (e.g., supervisor agent silently stops; agent in a "
        "reasoning loop) — how do you detect and recover from these in production? "
        "What automated evaluation frameworks have worked for detecting regressions after a system prompt or model update?"
    )
    _body(sl, questions, Inches(0.75), height=Inches(5.9), font_pt=10.5)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    print("Opening source template: %s" % SRC)
    prs = Presentation(SRC)

    print("Removing all existing slides...")
    _remove_all_slides(prs)

    print("Building Operations slides...")
    slide_cover(prs)
    slide_playing_rules(prs, 2)
    slide_agenda(prs)
    slide_platform_arch(prs)
    slide_op_principles(prs)
    # Release Management
    slide_agent_registry(prs)
    slide_deployment_lifecycle(prs)
    slide_testing_strategy(prs)
    # Observability
    slide_evalops(prs)
    slide_observability(prs)
    slide_behavioural_audit(prs)
    # Incident Management
    slide_hitl(prs)
    slide_emergency_stop(prs)
    # Platform Reliability & SLAs
    slide_reliability(prs)
    slide_sla(prs)
    # FinOps
    slide_finops(prs)
    slide_cost_comparison(prs)
    # Feedback & Lifecycle
    slide_feedback_loop(prs)
    slide_model_lifecycle(prs)
    # Operating Model
    slide_raci(prs)
    slide_ops_model(prs)
    slide_runbooks(prs)
    # Open Questions
    slide_open_questions(prs)

    print("Saving to: %s" % DST)
    prs.save(DST)
    print("Repacking zip to remove stale template slides...")
    _repack_zip(DST)
    print("Done. Total slides: %d" % len(prs.slides))


if __name__ == "__main__":
    main()
